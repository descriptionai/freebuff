# -*- coding: utf-8 -*-
"""
발표자료/프리젠테이션.txt 조건에 따라
은하 우주 배경 + 태양계 컨셉 발표 PPT(.pptx)를 생성한다.

- 배경: 첨부 이미지(나선 은하)를 PIL로 재현한 은하 우주 이미지를 매 슬라이드 배경으로 사용
  (진한 검정 우주 + 나선 은하(우상→좌하 대각선) + 밝은 코어 + 별 + 위성 은하)
- 인트로: 경고 로고 상단 중앙, 하단에 '시간허비 · 긁어 부스럼 · 업무부담 발생' 3문구만
- 메인: 태양(성공)을 중심으로 8개 행성이 궤도 위 다양한 위치에 배치
  (행성 클릭 -> 상세 슬라이드(정중앙 80%), 오른쪽 상단 X 버튼 -> 메인 복귀)
- 상세 슬라이드 9장: 성공/학습/데이터/화면/자동/교육/회의/결과/검증
- 마지막: 태양계를 2차원 원근(맨 마지막 행성 뒤)에서 바라본 슬라이드
  (태양은 정중앙 밝은 점, 원근법으로 행성들이 작아짐, 모두 화면에 보임)

※ PPT 정적 파일 특성상 '행성에서 뿜어져 나오는/빨려 들어가는' 애니메이션과
   복귀 후 행성 테두리 표시는 슬라이드 전환·하이퍼링크로 대체된다.
"""
import math
import os
import random

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

SW, SH = 13.333, 7.5
CX, CY = SW / 2.0, SH / 2.0
FONT = "맑은 고딕"

C_PANEL = RGBColor(0x0F, 0x1D, 0x4A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LBLUE = RGBColor(0x9F, 0xE0, 0xFF)
C_DIM = RGBColor(0x9A, 0xAA, 0xCF)
C_RED = RGBColor(0xE8, 0x4A, 0x4A)
C_GREEN = RGBColor(0x5A, 0xE8, 0x9A)
C_YELLOW = RGBColor(0xFF, 0xD9, 0x4A)

BG_IMG = os.path.join("발표자료", "은하_배경.png")
IMG_DIR = os.path.join("발표자료", "이미지")

# --------------------------------------------------------------------------
# 1. 은하 배경 이미지 생성 (PIL)
# --------------------------------------------------------------------------

def _ellipse_layer(w, h, color, alpha, blur=8):
    layer = Image.new("RGBA", (int(w) + blur * 4 + 8, int(h) + blur * 4 + 8), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)
    d.ellipse([0, 0, w, h], fill=color + (alpha,))
    if blur:
        layer = layer.filter(ImageFilter.GaussianBlur(blur))
    return layer


def _paste_rotated(base, layer, cx, cy, angle_deg):
    rot = layer.rotate(angle_deg, resample=Image.BICUBIC, expand=True)
    w, h = rot.size
    base.alpha_composite(rot, (int(cx - w / 2), int(cy - h / 2)))


def gen_galaxy_bg(path, seed=42):
    """첨부 이미지(나선 은하)를 본뜬 1920x1080 은하 우주 배경을 생성한다."""
    random.seed(seed)
    W, H = 1920, 1080

    # 기본: 거의 검정, 위쪽에 은은한 남색 그라데이션
    img = Image.new("RGB", (W, H), (2, 3, 8))
    d = ImageDraw.Draw(img)
    for y in range(H):
        t = y / H
        r = int(10 + (2 - 10) * t)
        g = int(18 + (3 - 18) * t)
        b = int(46 + (8 - 46) * t)
        d.line([(0, y), (W, y)], fill=(r, g, b))

    gx, gy = 1400, 300        # 나선 은하 중심 (우상)
    ang = -38                 # 대각선 방향 (우상 -> 좌하)
    gal = Image.new("RGBA", (W, H), (0, 0, 0, 0))

    # 은하 원반: 겹겹 타원 (바깥 헤일로 -> 안쪽 원반)
    _paste_rotated(gal, _ellipse_layer(1150, 420, (150, 170, 235), 26, blur=30), gx, gy, ang)
    _paste_rotated(gal, _ellipse_layer(1030, 330, (95, 115, 200), 55, blur=14), gx, gy, ang)
    _paste_rotated(gal, _ellipse_layer(920, 285, (135, 150, 225), 70, blur=10), gx, gy, ang)
    _paste_rotated(gal, _ellipse_layer(700, 215, (175, 185, 240), 85, blur=8), gx, gy, ang)

    # 먼지 띠 (어두운 갈색)
    _paste_rotated(gal, _ellipse_layer(880, 60, (52, 38, 26), 130, blur=10), gx - 60, gy - 70, ang + 4)
    _paste_rotated(gal, _ellipse_layer(780, 45, (45, 32, 22), 120, blur=10), gx + 90, gy + 55, ang - 3)

    # 중심부 밝은 코어 (은은한 노랑 -> 흰색)
    _paste_rotated(gal, _ellipse_layer(430, 150, (255, 224, 170), 120, blur=18), gx, gy, ang)
    _paste_rotated(gal, _ellipse_layer(230, 82, (255, 240, 200), 170, blur=8), gx, gy, ang)
    _paste_rotated(gal, _ellipse_layer(90, 34, (255, 255, 240), 230, blur=3), gx, gy, ang)

    # 나선 팔: 별무리(별 형성 영역) 블롭을 나선 곡선 위에 배치
    squash = 0.36
    ca, sa = math.cos(math.radians(ang)), math.sin(math.radians(ang))
    r0, r1, turns = 40, 560, 2.1
    k = (r1 - r0) / (turns * 2 * math.pi)
    palette = [(190, 210, 255), (230, 235, 255), (170, 185, 255),
               (255, 195, 215), (195, 175, 255), (215, 235, 255)]
    for arm in range(2):
        t = 0.0
        while True:
            r = r0 + k * t
            if r > r1:
                break
            theta = t + arm * math.pi
            u = r * math.cos(theta)
            v = r * math.sin(theta) * squash
            x = gx + u * ca - v * sa
            y = gy + u * sa + v * ca
            if 0 < x < W and 0 < y < H and random.random() < 0.6:
                col = random.choice(palette)
                size = random.uniform(7, 24)
                blob = _ellipse_layer(size, size, col, random.randint(70, 130), blur=4)
                gal.alpha_composite(blob, (int(x - size / 2), int(y - size / 2)))
            t += 16.0 / math.sqrt(k * k + r * r)   # 호 길이 기준 일정 간격

    gal = gal.filter(ImageFilter.GaussianBlur(1.5))

    # 위성 은하 (우하, M32 느낌)
    sat = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    _paste_rotated(sat, _ellipse_layer(150, 100, (215, 225, 255), 40, blur=16), 1560, 770, -20)
    _paste_rotated(sat, _ellipse_layer(70, 45, (235, 240, 255), 90, blur=6), 1560, 770, -20)
    _paste_rotated(sat, _ellipse_layer(26, 16, (255, 255, 255), 200, blur=2), 1560, 770, -20)

    # 별들 (다수 + 밝은 별 몇 개는 십자 광선)
    stars = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    sd = ImageDraw.Draw(stars)
    for _ in range(1700):
        x = random.uniform(0, W)
        y = random.uniform(0, H)
        r = random.choice([1, 1, 1, 2, 2, 3])
        br = random.randint(110, 255)
        col = (br, br, min(255, br + 25))
        sd.ellipse([x - r, y - r, x + r, y + r], fill=col + (255,))
    for _ in range(28):  # 밝은 별: 십자 광선
        x = random.uniform(40, W - 40)
        y = random.uniform(40, H - 40)
        br = random.randint(200, 255)
        col = (br, br, min(255, br + 20))
        ln = random.uniform(10, 22)
        sd.line([x - ln, y, x + ln, y], fill=col + (200,))
        sd.line([x, y - ln, x, y + ln], fill=col + (200,))
        sd.ellipse([x - 2, y - 2, x + 2, y + 2], fill=col + (255,))

    out = Image.alpha_composite(img.convert("RGBA"), gal)
    out = Image.alpha_composite(out, sat)
    out = Image.alpha_composite(out, stars)
    out.convert("RGB").save(path)
    print("GENERATED BG:", path)


# --------------------------------------------------------------------------
# 2. PPT 제작
# --------------------------------------------------------------------------

if not os.path.exists(BG_IMG):
    gen_galaxy_bg(BG_IMG)

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
blank = prs.slide_layouts[6]


def _shape(slide, kind, x, y, w, h, fill, line=None, lw=1.0, dash=None, rot=None):
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
        if dash:
            shp.line.dash_style = dash
    if rot is not None:
        shp.rotation = rot
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, s, size=14, color=C_WHITE, bold=False,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT, spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': font})
    rPr.append(ea)
    return tb


def add_bg(slide):
    """은하 배경 이미지를 슬라이드 전체에 깐다 (첫 번째 셰이프 = 맨 아래)."""
    slide.shapes.add_picture(BG_IMG, 0, 0, width=prs.slide_width, height=prs.slide_height)


def add_slide():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    return slide


def link_to(shape, slide, target):
    """셰이프를 클릭하면 target 슬라이드로 이동하는 내부 하이퍼링크."""
    rId = slide.part.relate_to(target.part, RT.SLIDE, is_external=False)
    sp = shape._element
    cNvPr = sp.nvSpPr.cNvPr
    hlink = cNvPr.makeelement(qn('a:hlinkClick'), {})
    hlink.set(qn('r:id'), rId)
    hlink.set('action', 'ppaction://hlinksldjump')
    cNvPr.append(hlink)


def check_box(slide, x, y, size=0.2):
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, size, size,
           RGBColor(0x1E, 0x3A, 0x6E), line=C_LBLUE, lw=1.2)
    _text(slide, x, y - 0.012, size, size, "✓", size=11, color=C_GREEN, bold=True)


def checkbox_list(slide, x, y, items, size=13, gap=0.44, color=C_WHITE):
    for i, item in enumerate(items):
        yy = y + i * gap
        check_box(slide, x, yy)
        _text(slide, x + 0.3, yy - 0.03, 6.8, gap - 0.05, item, size=size,
              color=color, align=PP_ALIGN.LEFT)


def image_placeholder(slide, x, y, w, h, caption):
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
           RGBColor(0x0C, 0x16, 0x38), line=RGBColor(0x4A, 0x6A, 0x9A), lw=1.2,
           dash=MSO_LINE.DASH)
    _text(slide, x, y + h * 0.32, w, 0.38, "이미지 삽입 위치", size=11, color=C_DIM)
    _text(slide, x + 0.15, y + h * 0.52, w - 0.3, h * 0.42, caption,
          size=10.5, color=C_LBLUE, spacing=1.05)


def put_image(slide, x, y, w, h, fname, caption):
    """발표자료/이미지의 삽화가 있으면 그 이미지를, 없으면 플레이스홀더를 넣는다."""
    path = os.path.join(IMG_DIR, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        image_placeholder(slide, x, y, w, h, caption)


def browser(slide, x, y, w, h, title):
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
           RGBColor(0x16, 0x22, 0x46), line=RGBColor(0x5A, 0x7A, 0xAA), lw=1.2)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.08, y + 0.1, w - 0.16, 0.28,
           RGBColor(0x2A, 0x3A, 0x62))
    _text(slide, x + 0.18, y + 0.1, w - 0.36, 0.28, title, size=9,
          color=C_LBLUE, align=PP_ALIGN.LEFT)
    return (x + 0.15, y + 0.5, w - 0.3, h - 0.65)


def truck(slide, x, y, w, label, ccolor):
    h = 0.62
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w * 0.62, h, ccolor, line=C_WHITE, lw=1.0)
    _text(slide, x, y - 0.03, w * 0.62, h + 0.06, label, size=13, bold=True)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + w * 0.62 + 0.03, y + 0.12,
           w * 0.30, h * 0.72, RGBColor(0x4A, 0x6A, 0x8A), line=C_WHITE, lw=1.0)
    for wx in (x + 0.1, x + w * 0.46, x + w * 0.72, x + w * 0.88):
        _shape(slide, MSO_SHAPE.OVAL, wx, y + h - 0.1, 0.15, 0.15,
               RGBColor(0x11, 0x11, 0x11), line=C_DIM, lw=0.5)


# ---------- 1. 인트로 (경고 로고 + 하단 3문구만) ----------
s_intro = add_slide()
# 경고 로고: 상단 중앙 (빨간 원 테두리 + 노란 삼각형 + !)
_shape(s_intro, MSO_SHAPE.OVAL, CX - 1.25, 0.75, 2.5, 2.5, None,
       line=C_RED, lw=3.0)
_shape(s_intro, MSO_SHAPE.OVAL, CX - 1.12, 0.88, 2.24, 2.24, None,
       line=RGBColor(0x8A, 0x2A, 0x2A), lw=1.0)
_shape(s_intro, MSO_SHAPE.ISOSCELES_TRIANGLE, CX - 0.95, 1.2, 1.9, 1.65, C_YELLOW)
_text(s_intro, CX - 0.95, 1.3, 1.9, 1.45, "!", size=52, bold=True,
      color=RGBColor(0x33, 0x26, 0x00))
# 하단 3문구
intro_words = ["시간허비", "긁어 부스럼", "업무부담 발생"]
for i, w in enumerate(intro_words):
    wx = 1.7 + i * 3.31
    _shape(s_intro, MSO_SHAPE.ROUNDED_RECTANGLE, wx, 6.0, 2.6, 0.62,
           RGBColor(0x2A, 0x0E, 0x0E), line=C_RED, lw=1.2)
    _text(s_intro, wx, 6.02, 2.6, 0.58, w, size=19, bold=True, color=C_RED)

# ---------- 2. 메인 (태양계) ----------
s_main = add_slide()
_text(s_main, 0, 0.28, SW, 0.6, "우주 프로젝트 — 성공을 향한 태양계", size=26, bold=True, color=C_YELLOW)
_text(s_main, 0, 0.92, SW, 0.4,
      "태양(성공)과 8개 행성 · 각 행성이나 태양을 클릭하면 상세 내용이 나옵니다",
      size=13, color=C_LBLUE)

# 태양 (성공)
sun_r = 0.55
for sr, sc in ((1.05, RGBColor(0x8A, 0x6A, 0x14)), (0.85, RGBColor(0xC8, 0x9A, 0x2A)),
               (0.68, RGBColor(0xFF, 0xC8, 0x3A)), (sun_r, C_YELLOW)):
    _shape(s_main, MSO_SHAPE.OVAL, CX - sr, CY - sr, sr * 2, sr * 2, sc)
sun_shape = None
for shp in s_main.shapes:
    if shp.width is not None and shp.width.inches == sun_r * 2:
        sun_shape = shp
_text(s_main, CX - 1.2, CY + sun_r + 0.06, 2.4, 0.35, "성공 (태양)", size=13, bold=True, color=C_YELLOW)

# 행성 정의: (이름, 색, 궤도 반지름, 각도, 행성 반지름)
planets = [
    ("학습", RGBColor(0x9E, 0x9E, 0x9E), 1.30, 205, 0.24),
    ("데이터", RGBColor(0xE8, 0xC5, 0x4A), 1.72, 155, 0.25),
    ("화면", RGBColor(0x4A, 0x9B, 0xD8), 2.14, 212, 0.26),
    ("자동", RGBColor(0xC0, 0x4A, 0x3A), 2.56, 148, 0.27),
    ("교육", RGBColor(0xD9, 0x8A, 0x4A), 2.98, 218, 0.28),
    ("회의", RGBColor(0xE8, 0xD8, 0x8A), 3.40, 142, 0.29),
    ("결과", RGBColor(0x8A, 0xD8, 0xE8), 3.82, 205, 0.30),
    ("검증", RGBColor(0x4A, 0x6A, 0xC8), 4.24, 172, 0.31),
]

planet_shapes = []
for name, color, rx, deg, pr in planets:
    ry = rx * 0.52
    _shape(s_main, MSO_SHAPE.OVAL, CX - rx, CY - ry, rx * 2, ry * 2, None,
           line=RGBColor(0x5A, 0x7A, 0xAA), lw=0.75)
    rad = math.radians(deg)
    px = CX + rx * math.cos(rad)
    py = CY + ry * math.sin(rad)
    planet = _shape(s_main, MSO_SHAPE.OVAL, px - pr, py - pr, pr * 2, pr * 2, color,
                    line=C_WHITE, lw=1.0)
    planet_shapes.append(planet)
    _text(s_main, px - 1.1, py + pr + 0.02, 2.2, 0.3, name, size=11.5, bold=True)
    if name == "회의":
        _shape(s_main, MSO_SHAPE.OVAL, px - pr - 0.14, py - pr * 0.55, pr * 2 + 0.28,
               pr * 1.1, None, line=RGBColor(0xE8, 0xD8, 0x8A), lw=1.2, rot=-18)

# ---------- 상세 슬라이드 공통 ----------
def detail_slide(title):
    slide = add_slide()
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.33, 0.75, 10.67, 6.0, C_PANEL,
           line=RGBColor(0x3A, 0x5A, 0x9A), lw=1.5)
    _text(slide, 1.33, 0.95, 10.0, 0.55, title, size=25, bold=True, color=C_YELLOW,
          align=PP_ALIGN.LEFT)
    xb = _shape(slide, MSO_SHAPE.OVAL, 11.62, 1.02, 0.42, 0.42, C_RED, line=C_WHITE, lw=1.2)
    _text(slide, 11.62, 1.05, 0.42, 0.36, "✕", size=14, bold=True)
    link_to(xb, slide, s_main)
    return slide

# ---------- 3. 성공 ----------
s_success = detail_slide("성 공 · 실 패 — 두 개의 차선")
for lane_y, words, tsize in (
        (2.0, "책임 · 해보자 · 다각화 · 변화 · 오픈마인드 · 긍정적 태도 · 칭찬 · 열정 · 희생", 11),
        (4.2, "회피 · 무책임 · 귀찮음 · 숨기기 · 어려움 · 몰아주기 · 떠넘기기 · 억제 · 무시 · 단정 · 부정적 초점 · 기존틀 고수 · 열정식히기", 9.5)):
    _shape(s_success, MSO_SHAPE.RECTANGLE, 1.6, lane_y, 10.1, 1.5, RGBColor(0x2E, 0x2E, 0x38),
           line=RGBColor(0x6A, 0x6A, 0x76), lw=1.0)
    for i in range(6):
        _shape(s_success, MSO_SHAPE.RECTANGLE, 1.9 + i * 1.6, lane_y + 0.72, 0.2, 0.06, C_WHITE)
    _text(s_success, 4.6, lane_y + 0.32, 6.8, 0.9, words, size=tsize, color=C_WHITE, spacing=1.1)
truck(s_success, 1.7, 2.38, 2.6, "성공", RGBColor(0x2A, 0x7A, 0x4A))
truck(s_success, 1.7, 4.58, 2.6, "실패", RGBColor(0x8A, 0x2A, 0x2A))

# ---------- 4. 학습 ----------
s_learn = detail_slide("학 습 — 배움과 훈련")
put_image(s_learn, 1.9, 1.95, 4.6, 3.85, "학습_스파르타.png",
          "스파르타 군인들이 컴퓨터를 둘러싸고 열정적으로 가르치는 모습")
checkbox_list(s_learn, 7.0, 1.95, [
    "용어 정리 및 통일", "시즌 분류 기준 설정", "결과 도출 방법론 연구",
    "작은 단위 목표 수립", "단계별 첵크 리스트", "설명 및 테스트 방식",
    "AI 상향 평준화", "최종 검증 때 사용",
])

# ---------- 5. 데이터 ----------
s_data = detail_slide("데 이 터 — 추출과 가공")
_shape(s_data, MSO_SHAPE.RECTANGLE, 2.1, 3.35, 2.3, 1.5, RGBColor(0x3A, 0x44, 0x54),
       line=RGBColor(0x8A, 0x9A, 0xAA), lw=1.0)
_shape(s_data, MSO_SHAPE.RECTANGLE, 2.45, 3.05, 1.6, 0.32, RGBColor(0x4A, 0x56, 0x66))
for i in range(3):
    _shape(s_data, MSO_SHAPE.RECTANGLE, 2.35, 4.78 + i * 0.22, 1.5, 0.2, C_WHITE)
_text(s_data, 1.8, 5.55, 2.9, 0.3, "종이 데이터", size=10.5, color=C_LBLUE)
_shape(s_data, MSO_SHAPE.RECTANGLE, 4.6, 3.0, 1.25, 1.25, C_RED, line=C_WHITE, lw=1.0)
_shape(s_data, MSO_SHAPE.RECTANGLE, 5.07, 2.95, 0.3, 1.35, RGBColor(0xE8, 0xD8, 0x8A))
_shape(s_data, MSO_SHAPE.RECTANGLE, 4.55, 3.47, 1.35, 0.3, RGBColor(0xE8, 0xD8, 0x8A))
_text(s_data, 3.4, 4.45, 3.6, 0.3, "빨간색 포장 선물", size=10.5, color=C_LBLUE)
_text(s_data, 1.9, 2.0, 4.6, 0.35, "프린터에서 출력되는 데이터", size=12, color=C_WHITE, bold=True)
checkbox_list(s_data, 7.0, 1.95, [
    "데이터 추출 경로 정리", "필수 데이터로 가공", "시즌 분류 및 변화율 측정",
    "자동화 데이터 추출 방법", "비정형 통계 요청안 정의", "작은 단위 목표 수립",
    "단계별 첵크 리스트", "포장 데이터 추출",
])

# ---------- 6. 화면 ----------
s_screen = detail_slide("화 면 — 세 단계 구성")
ax, ay, aw, ah = browser(s_screen, 1.7, 2.0, 3.0, 3.6, "지도 화면")
for gx in (0, 1, 2, 3):
    _shape(s_screen, MSO_SHAPE.RECTANGLE, ax, ay + gx * (ah / 4), aw, 0.015,
           RGBColor(0x5A, 0x7A, 0xAA))
for gy in (0, 1, 2, 3):
    _shape(s_screen, MSO_SHAPE.RECTANGLE, ax + gy * (aw / 4), ay, 0.015, ah,
           RGBColor(0x5A, 0x7A, 0xAA))
_text(s_screen, ax, ay + ah * 0.4, aw, 0.35, "지도", size=12, color=C_LBLUE)
mid_x, mid_y, mid_w = 5.2, 2.1, 2.9
_shape(s_screen, MSO_SHAPE.HEXAGON, mid_x + 0.95, mid_y, 1.0, 0.9, RGBColor(0x8A, 0x9A, 0xAA))
_text(s_screen, mid_x + 0.95, mid_y + 0.95, 1.0, 0.3, "나사", size=11, color=C_WHITE)
_shape(s_screen, MSO_SHAPE.OVAL, mid_x + 1.05, mid_y + 1.35, 0.8, 0.8, RGBColor(0xC8, 0x8A, 0x4A))
_text(s_screen, mid_x + 0.95, mid_y + 2.2, 1.0, 0.3, "몽키스피너", size=9.5, color=C_WHITE)
_shape(s_screen, MSO_SHAPE.OVAL, mid_x + 1.05, mid_y + 2.6, 0.8, 0.8, RGBColor(0x8A, 0xA8, 0xC8))
for gd in (0, 45, 90, 135, 180, 225, 270, 315):
    _shape(s_screen, MSO_SHAPE.RECTANGLE,
           mid_x + 1.32, mid_y + 2.46 + 0.28 * math.sin(math.radians(gd)),
           0.26, 0.26, RGBColor(0x8A, 0xA8, 0xC8))
_text(s_screen, mid_x + 0.85, mid_y + 3.45, 1.2, 0.3, "톱니바퀴", size=10, color=C_WHITE)
_shape(s_screen, MSO_SHAPE.OVAL, 8.9, 2.7, 2.5, 1.5, C_WHITE)
_shape(s_screen, MSO_SHAPE.OVAL, 9.75, 3.05, 0.8, 0.8, RGBColor(0x4A, 0x9B, 0xD8))
_shape(s_screen, MSO_SHAPE.OVAL, 10.0, 3.3, 0.3, 0.3, RGBColor(0x11, 0x11, 0x11))
_text(s_screen, 8.9, 4.3, 2.5, 0.35, "눈 (1개, 크게)", size=12, bold=True)
_shape(s_screen, MSO_SHAPE.RIGHT_ARROW, 4.72, 3.35, 0.5, 0.4, C_LBLUE)
_shape(s_screen, MSO_SHAPE.RIGHT_ARROW, 8.12, 3.35, 0.5, 0.4, C_LBLUE)

# ---------- 7. 자동 ----------
s_auto = detail_slide("자 동 — 순환 프로세스")
put_image(s_auto, 4.9, 2.55, 3.5, 2.7, "자동_무써기.png",
          "GIF: 도마 위에서 무를 칼로 채 써는 모습")
_shape(s_auto, MSO_SHAPE.OVAL, CX - 2.6, CY - 1.7, 5.2, 3.4, None,
       line=RGBColor(0x5A, 0x7A, 0xAA), lw=1.2, dash=MSO_LINE.DASH)
arrows = [
    (MSO_SHAPE.RIGHT_ARROW, CX - 0.3, 1.35, 0.6, 0.4, 0),
    (MSO_SHAPE.RIGHT_ARROW, 9.75, CY - 0.2, 0.6, 0.4, 90),
    (MSO_SHAPE.RIGHT_ARROW, CX - 0.3, 5.65, 0.6, 0.4, 180),
    (MSO_SHAPE.RIGHT_ARROW, 2.95, CY - 0.2, 0.6, 0.4, 270),
]
for kind, x, y, w, h, rot in arrows:
    _shape(s_auto, kind, x, y, w, h, C_LBLUE, rot=rot)
_text(s_auto, CX - 0.8, 0.9, 1.6, 0.4, "문 제", size=14, bold=True, color=C_RED)
_text(s_auto, 9.95, CY - 0.5, 1.8, 0.4, "AI", size=14, bold=True, color=C_LBLUE)
_text(s_auto, CX - 0.8, 6.25, 1.6, 0.4, "학 습", size=14, bold=True, color=C_GREEN)
_text(s_auto, 1.5, CY - 0.5, 1.8, 0.4, "결 과", size=14, bold=True, color=C_YELLOW)

# ---------- 8. 교육 ----------
s_edu = detail_slide("교 육 — 지식의 흐름")
image_placeholder(s_edu, 2.0, 1.8, 3.3, 1.9, "자료 이미지")
image_placeholder(s_edu, 8.0, 1.8, 3.3, 1.9, "화면 이미지")
_shape(s_edu, MSO_SHAPE.OVAL, CX - 0.3, 3.3, 0.6, 0.6, RGBColor(0xE8, 0xC8, 0xA0))
_shape(s_edu, MSO_SHAPE.ROUNDED_RECTANGLE, CX - 0.45, 3.95, 0.9, 1.0, RGBColor(0x4A, 0x6A, 0x8A))
_text(s_edu, CX - 1.4, 2.4, 2.8, 0.4, "지식이 머리로 흡수", size=13, bold=True, color=C_YELLOW)
_shape(s_edu, MSO_SHAPE.DOWN_ARROW, CX - 0.22, 2.85, 0.44, 0.5, C_LBLUE)
_text(s_edu, CX - 1.4, 5.05, 2.8, 0.35, "사람", size=12, color=C_LBLUE)
truck(s_edu, 1.8, 4.5, 2.5, "A to Z", RGBColor(0x2A, 0x6A, 0x9A))
_text(s_edu, 1.8, 5.25, 2.5, 0.3, "컨테이너 트럭 (A to Z)", size=10, color=C_LBLUE)

# ---------- 9. 회의 ----------
s_meet = detail_slide("회 의 — 원탁 논의")
table_cx, table_cy = 3.9, 4.1
_shape(s_meet, MSO_SHAPE.OVAL, table_cx - 1.85, table_cy - 1.05, 3.7, 2.1,
       RGBColor(0x6A, 0x4A, 0x2A), line=RGBColor(0xA8, 0x7A, 0x4A), lw=1.2)
papers = ["책임", "해보자", "다각화", "변화", "오픈마인드", "긍정적 태도", "칭찬", "열정", "희생"]
for i, pw in enumerate(papers):
    gx, gy = i % 3, i // 3
    px = table_cx - 1.32 + gx * 0.92
    py = table_cy - 0.42 + gy * 0.36
    _shape(s_meet, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, 0.85, 0.3, C_WHITE)
    _text(s_meet, px, py - 0.02, 0.85, 0.32, pw, size=7.5, bold=True,
          color=RGBColor(0x33, 0x33, 0x33))
people = [(30, "여"), (90, "남"), (150, "여"), (210, "남"), (270, "여"), (330, "남")]
for deg, g in people:
    rad = math.radians(deg)
    px = table_cx + 2.15 * math.cos(rad)
    py = table_cy + 1.25 * math.sin(rad)
    _shape(s_meet, MSO_SHAPE.OVAL, px - 0.2, py - 0.2, 0.4, 0.4,
           RGBColor(0xE8, 0xC8, 0xA0), line=C_WHITE, lw=0.8)
    _text(s_meet, px - 0.2, py - 0.42, 0.4, 0.22, g, size=8, bold=True)
_text(s_meet, 1.9, 1.85, 4.0, 0.35, "원형 탁자 — 여성 3명 · 남성 3명", size=12,
      color=C_WHITE, bold=True)
checkbox_list(s_meet, 7.0, 1.95, [
    "진행현황 공유", "조건 설정 및 검증", "절차 추가/개선 사항도출",
    "새로운 아이디어", "질의 응답",
], gap=0.55)

# ---------- 10. 결과 ----------
s_result = detail_slide("결 과 — 성과의 두 얼굴")
_shape(s_result, MSO_SHAPE.ROUNDED_RECTANGLE, 1.7, 1.75, 4.9, 4.3, RGBColor(0x0E, 0x2A, 0x1E),
       line=C_GREEN, lw=1.4)
_text(s_result, 1.7, 1.9, 4.9, 0.4, "만족 성과", size=16, bold=True, color=C_GREEN)
_shape(s_result, MSO_SHAPE.RECTANGLE, 2.3, 3.15, 1.5, 1.4, RGBColor(0xE8, 0xE8, 0xE0),
       line=RGBColor(0x8A, 0x8A, 0x80), lw=1.0)
_shape(s_result, MSO_SHAPE.ISOSCELES_TRIANGLE, 2.22, 2.85, 1.66, 0.5, RGBColor(0x8A, 0x4A, 0x3A))
_shape(s_result, MSO_SHAPE.RECTANGLE, 2.85, 3.95, 0.4, 0.6, RGBColor(0x5A, 0x3A, 0x2A))
_shape(s_result, MSO_SHAPE.OVAL, 2.6, 3.4, 0.3, 0.3, RGBColor(0x2A, 0x6A, 0x9A))
_text(s_result, 2.3, 4.6, 1.5, 0.3, "우체국", size=10, bold=True, color=RGBColor(0x33, 0x33, 0x33))
put_image(s_result, 2.1, 5.0, 4.1, 0.9, "결과_우체국행진.png",
          "정장 여성 3명 · 남성 3명이 우체국으로 걸어가는 모습")
_shape(s_result, MSO_SHAPE.ROUNDED_RECTANGLE, 6.75, 1.75, 4.9, 4.3, RGBColor(0x2E, 0x14, 0x14),
       line=C_RED, lw=1.4)
_text(s_result, 6.75, 1.9, 4.9, 0.4, "불만족 성과", size=16, bold=True, color=C_RED)
_shape(s_result, MSO_SHAPE.RECTANGLE, 8.35, 3.5, 1.6, 1.6, C_RED, line=C_WHITE, lw=1.0)
_shape(s_result, MSO_SHAPE.RECTANGLE, 9.05, 3.4, 0.2, 1.8, RGBColor(0xE8, 0xD8, 0x8A))
_shape(s_result, MSO_SHAPE.RECTANGLE, 8.25, 4.2, 1.8, 0.2, RGBColor(0xE8, 0xD8, 0x8A))
_shape(s_result, MSO_SHAPE.OVAL, 8.6, 2.75, 1.1, 1.1, RGBColor(0x8A, 0x6A, 0x4A),
       line=C_WHITE, lw=1.0)
for hx, hr in ((8.62, 0.5), (9.68, 0.5)):
    _shape(s_result, MSO_SHAPE.ISOSCELES_TRIANGLE, hx, 2.5, 0.28, 0.42,
           RGBColor(0xE8, 0xE8, 0xE0), rot=hr)
_shape(s_result, MSO_SHAPE.OVAL, 8.95, 3.0, 0.16, 0.16, C_WHITE)
_shape(s_result, MSO_SHAPE.OVAL, 9.35, 3.0, 0.16, 0.16, C_WHITE)
_shape(s_result, MSO_SHAPE.OVAL, 9.05, 3.62, 0.36, 0.2, None, line=C_RED, lw=1.5)
_text(s_result, 8.35, 5.2, 1.6, 0.3, "코뚜레한 소 머리", size=10, color=C_LBLUE)

# ---------- 11. 검증 ----------
s_check = detail_slide("검 증 — 다각도 확인")
put_image(s_check, 4.85, 2.6, 3.6, 2.7, "검증_종이날림.png",
          "GIF: 컴퓨터 앞에서 종이를 작업하다 하나씩 날리는 모습")
ax, ay, aw, ah = browser(s_check, 1.7, 1.55, 3.1, 2.1, "네이버 맵")
for gx in (0, 1, 2):
    _shape(s_check, MSO_SHAPE.RECTANGLE, ax, ay + gx * (ah / 3), aw, 0.012, RGBColor(0x5A, 0x7A, 0xAA))
for gy in (0, 1, 2):
    _shape(s_check, MSO_SHAPE.RECTANGLE, ax + gy * (aw / 3), ay, 0.012, ah, RGBColor(0x5A, 0x7A, 0xAA))
_text(s_check, ax, ay + ah * 0.4, aw, 0.3, "지도", size=10, color=C_LBLUE)
ax, ay, aw, ah = browser(s_check, 8.5, 1.55, 3.1, 2.1, "제미나이")
_text(s_check, ax, ay + ah * 0.35, aw, 0.4, "Gemini", size=13, bold=True, color=C_LBLUE)
_shape(s_check, MSO_SHAPE.ROUNDED_RECTANGLE, 1.7, 4.5, 3.1, 2.0, RGBColor(0x22, 0x2A, 0x3A),
       line=RGBColor(0x5A, 0x7A, 0xAA), lw=1.2)
_shape(s_check, MSO_SHAPE.ROUNDED_RECTANGLE, 1.85, 4.65, 2.8, 0.45, RGBColor(0xC8, 0xE8, 0xC8))
for bi in range(9):
    bx = 1.85 + (bi % 3) * 1.0
    by = 5.25 + (bi // 3) * 0.42
    _shape(s_check, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, 0.85, 0.32, RGBColor(0x4A, 0x56, 0x66))
_text(s_check, 1.7, 6.35, 3.1, 0.28, "계산기", size=10, color=C_LBLUE)
_shape(s_check, MSO_SHAPE.ROUNDED_RECTANGLE, 8.5, 4.5, 3.1, 2.0, RGBColor(0x22, 0x2A, 0x3A),
       line=RGBColor(0x5A, 0x7A, 0xAA), lw=1.2)
_shape(s_check, MSO_SHAPE.RECTANGLE, 8.85, 4.8, 1.7, 1.3, C_WHITE, rot=-6)
_shape(s_check, MSO_SHAPE.RECTANGLE, 10.35, 4.85, 0.5, 0.09, RGBColor(0x2A, 0x4A, 0x8A), rot=30)
_text(s_check, 8.5, 6.35, 3.1, 0.28, "종이와 펜", size=10, color=C_LBLUE)

# ---------- 12. 마지막 (원근 태양계) ----------
s_end = add_slide()
_text(s_end, 0, 0.3, SW, 0.55, "끝에서 바라본 태양계 — 모두가 보입니다", size=24, bold=True,
      color=C_YELLOW)
for sr, sc in ((1.3, RGBColor(0x5A, 0x44, 0x0A)), (1.0, RGBColor(0x8A, 0x6A, 0x14)),
               (0.72, RGBColor(0xC8, 0x9A, 0x2A)), (0.45, RGBColor(0xFF, 0xD9, 0x4A)),
               (0.22, C_WHITE)):
    _shape(s_end, MSO_SHAPE.OVAL, CX - sr, CY - sr, sr * 2, sr * 2, sc)
for i, (name, color) in enumerate([
        ("학습", RGBColor(0x9E, 0x9E, 0x9E)), ("데이터", RGBColor(0xE8, 0xC5, 0x4A)),
        ("화면", RGBColor(0x4A, 0x9B, 0xD8)), ("자동", RGBColor(0xC0, 0x4A, 0x3A)),
        ("교육", RGBColor(0xD9, 0x8A, 0x4A)), ("회의", RGBColor(0xE8, 0xD8, 0x8A)),
        ("결과", RGBColor(0x8A, 0xD8, 0xE8)), ("검증", RGBColor(0x4A, 0x6A, 0xC8))]):
    rx = 1.15 + i * 0.44
    ry = rx * 0.5
    _shape(s_end, MSO_SHAPE.OVAL, CX - rx, CY - ry, rx * 2, ry * 2, None,
           line=RGBColor(0x4A, 0x66, 0x8A), lw=0.6)
    deg = 195 - i * 8
    rad = math.radians(deg)
    px = CX + rx * math.cos(rad)
    py = CY + ry * math.sin(rad)
    pr = max(0.05, 0.16 - i * 0.013)
    _shape(s_end, MSO_SHAPE.OVAL, px - pr, py - pr, pr * 2, pr * 2, color, line=C_DIM, lw=0.5)
    _text(s_end, px - 0.6, py + pr + 0.02, 1.2, 0.24, name, size=8.5, color=C_DIM)
_text(s_end, 0, 6.75, SW, 0.4, "태양은 정중앙에서 밝게 빛나며, 모든 행성이 화면에 보입니다",
      size=12, color=C_DIM)

# ---------- 링크 연결 ----------
if sun_shape is not None:
    link_to(sun_shape, s_main, s_success)
for planet, target in zip(planet_shapes, [s_learn, s_data, s_screen, s_auto, s_edu,
                                          s_meet, s_result, s_check]):
    link_to(planet, s_main, target)

out_path = "발표자료/우주_태양계_프로젝트_은하.pptx"
prs.save(out_path)
print("SAVED:", out_path)
print("SLIDES:", len(prs.slides._sldIdLst))
